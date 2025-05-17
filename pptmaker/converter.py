import base64
import io
import os
from dataclasses import dataclass, field
from typing import List, Dict, Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import yaml

try:
    import openai
except ImportError:  # pragma: no cover
    openai = None

@dataclass
class SlideImage:
    description: str
    filename: str | None = None
    data: str | None = None  # base64 encoded image data

@dataclass
class SlideTable:
    markdown: str

@dataclass
class SlideContent:
    title: str | None
    texts: List[str] = field(default_factory=list)
    images: List[SlideImage] = field(default_factory=list)
    tables: List[SlideTable] = field(default_factory=list)


def describe_image(image_bytes: bytes) -> str:
    """Use OpenAI to generate an alt text description for an image."""
    if openai is None:
        raise RuntimeError("openai package is required for image description")
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY environment variable not set")
    openai.api_key = api_key
    # Using GPT-4 vision or similar; here we just send a prompt with base64 image
    encoded = base64.b64encode(image_bytes).decode("ascii")
    prompt = (
        "Describe the following image from a PowerPoint slide in one or two sentences."
    )
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-vision-preview",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt,
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{encoded}",
                            },
                        },
                    ],
                }
            ],
            max_tokens=60,
        )
        return response.choices[0].message.content.strip()
    except Exception as exc:  # pragma: no cover - network related
        raise RuntimeError(f"Failed to get image description: {exc}")


def convert_table_to_markdown(table) -> str:
    """Convert a python-pptx table shape to markdown."""
    rows = []
    for row in table.rows:
        cells = [cell.text for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if not rows:
        return ""
    header = rows[0]
    separator = "| " + " | ".join(["---" for _ in table.columns]) + " |"
    return "\n".join([header, separator] + rows[1:])


def extract_slide_content(slide) -> SlideContent:
    title = None
    if slide.shapes.title:
        title = slide.shapes.title.text
    content = SlideContent(title=title)
    for shape in slide.shapes:
        if not shape.has_text_frame and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Extract image
            image_bytes = shape.image.blob
            description = describe_image(image_bytes)
            encoded = base64.b64encode(image_bytes).decode("ascii")
            content.images.append(
                SlideImage(description=description, data=encoded)
            )
        elif shape.has_text_frame:
            text = shape.text_frame.text
            if text:
                content.texts.append(text)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            md = convert_table_to_markdown(shape.table)
            content.tables.append(SlideTable(markdown=md))
    return content


def pptx_to_yaml(path: str) -> str:
    prs = Presentation(path)
    slides: List[Dict[str, Any]] = []
    for slide in prs.slides:
        content = extract_slide_content(slide)
        slides.append({
            "title": content.title,
            "texts": content.texts,
            "images": [image.__dict__ for image in content.images],
            "tables": [table.__dict__ for table in content.tables],
        })
    return yaml.safe_dump({"slides": slides}, sort_keys=False, allow_unicode=True)

