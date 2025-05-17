import yaml
from pptx import Presentation
from pptx.util import Inches
import pptmaker.converter as conv


def create_sample_pptx(path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Title"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "First bullet"
    p = tf.add_paragraph()
    p.text = "Second bullet"
    rows, cols = 2, 2
    left = top = Inches(1)
    width = height = Inches(2)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"
    prs.save(path)


def test_pptx_to_yaml(tmp_path, monkeypatch):
    pptx_path = tmp_path / "sample.pptx"
    create_sample_pptx(pptx_path)
    monkeypatch.setattr(conv, "describe_image", lambda x: "desc")
    yaml_text = conv.pptx_to_yaml(str(pptx_path))
    data = yaml.safe_load(yaml_text)
    slide = data["slides"][0]
    assert slide["title"] == "Title"
    assert "First bullet" in slide["texts"][0]
    assert "| H1 | H2 |" in slide["tables"][0]["markdown"]
